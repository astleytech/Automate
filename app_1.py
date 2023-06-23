import streamlit as st
import transformers
import docx
import pptx

# Load the pre-trained summarization model
model = transformers.pipeline("summarization")

# Load the uploaded document and extract its paragraphs
uploaded_file = st.file_uploader("Upload a file")
if uploaded_file is not None:
    doc = docx.Document(uploaded_file)
    paragraphs = [p.text for p in doc.paragraphs]

    # Generate summarized bullets for each paragraph using the summarization model
    bullets = []
    for paragraph in paragraphs:
        summary = model(paragraph, max_length=50, min_length=10, do_sample=False)[0]['summary_text'].strip()
        bullets.append(summary)

    # Create a PowerPoint file and populate it with the summarized bullets
    prs = pptx.Presentation()
    for bullet in bullets:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Summary"
        text_frame = slide.shapes.add_textbox(left=0, top=80, width=prs.slide_width, height=prs.slide_height - 80).text_frame
        p = text_frame.add_paragraph()
        p.text = bullet

    # Save the PowerPoint file and display a download link
    prs.save("output.pptx")
    st.download_button(label="Download Summary", data=open("output.pptx", "rb").read(), file_name="output.pptx")